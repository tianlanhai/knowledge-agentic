# -*- coding: utf-8 -*-
"""
上海宇羲伏天智能科技有限公司出品

文件级注释：PPTX 文档加载器
内部逻辑：使用 python-pptx 解析 PowerPoint 文件，提取文本内容
说明：轻量级替代方案，避免使用 unstructured（体积约 4-5GB）
"""

from typing import List
from langchain_core.documents import Document


class PPTXLoader:
    """
    类级注释：PowerPoint 文档加载器
    使用 python-pptx 库解析 .pptx 文件
    """

    def __init__(self, file_path: str):
        """
        函数级注释：初始化 PPTX 加载器
        参数：
            file_path: PPTX 文件路径
        """
        self.file_path = file_path

    def load(self) -> List[Document]:
        """
        函数级注释：加载 PPTX 文件并提取文本内容
        返回值：Document 对象列表
        """
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError(
                "python-pptx 包未安装。请运行: pip install python-pptx"
            )

        # 内部逻辑：打开 PPTX 文件
        prs = Presentation(self.file_path)

        # 内部变量：存储所有幻灯片的文本内容
        full_text = []

        # 内部逻辑：遍历所有幻灯片
        for slide_num, slide in enumerate(prs.slides):
            # 内部变量：当前幻灯片的文本
            slide_text = []

            # 内部逻辑：遍历幻灯片中的所有形状
            for shape in slide.shapes:
                # 内部逻辑：只处理包含文本的形状
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)

            # 内部逻辑：如果有文本，添加到总内容中
            if slide_text:
                full_text.append(f"[Slide {slide_num + 1}]\n" + "\n".join(slide_text))

        # 内部逻辑：创建 Document 对象
        content = "\n\n".join(full_text)

        metadata = {
            "source": self.file_path,
            "file_path": self.file_path,
        }

        return [Document(page_content=content, metadata=metadata)]
